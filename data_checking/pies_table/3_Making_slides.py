import glob
import pathlib
import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

def absoluteFilePaths(directory):
	path_list=[]
	for dirpath,_,filenames in os.walk(directory):
		for f in filenames:
			path_list.append(os.path.abspath(os.path.join(dirpath, f)))
	return path_list
def EG_slide(offset=0):
		thedir=scaffold_path
		EG_dirs=([ name for name in os.listdir(thedir) if os.path.isdir(os.path.join(thedir, name)) ])
		print('EG_dirs',EG_dirs)
		for egindex,EG in enumerate(EG_dirs):
			print('EG',EG)
		# Add text
			leftpin=offset
			toppin=egindex+1
			EG_path=scaffold_path+'/'+EG
		# Get one figure as example
			#path_to_find=EG_path
			example_reaction_path=absoluteFilePaths(EG_path)[0]
			example_reaction_fig = cluster_slide.shapes.add_picture(example_reaction_path, Inches(leftpin+0.2), Inches(toppin+0.2), height=Inches(1))
		# Add text to cluster's slide
			tl=Inches(leftpin)
			tt=Inches(toppin)
			tw=Inches(1)
			th=Inches(1)
			txBox_EG = cluster_slide.shapes.add_textbox(tl, tt, tw, th)
			tf_EG = txBox_EG.text_frame
			tf_EG.text = EG
			tl=Inches(leftpin+1)
			tt=Inches(toppin)
			tw=Inches(1)
			th=Inches(1)
			txBox_ref = cluster_slide.shapes.add_textbox(tl, tt, tw, th)
			tf_ref = txBox_ref.text_frame
#			print('example_reaction_path',example_reaction_path)
			part_of_paths=example_reaction_path[-25:].split('\\')
			ID_file_name=part_of_paths[-1]
			print('file_name',ID_file_name)
			tf_ref.text = ID_file_name.replace('.png','')
		# GO TO EACH OTHER LEAVING GROUP FOLDER
			thedir=EG_path
			EG_sub_dirs=([ name for name in os.listdir(thedir) if os.path.isdir(os.path.join(thedir, name)) ])
			print('EG_sub_dirs',EG_sub_dirs)
		# Create new slide-----------------------------------------
			EG_slide = prs.slides.add_slide(blank_slide_layout)\
		# Add text to cluster's slide
			tl=Inches(1)
			tt=Inches(0)
			tw=Inches(15)
			th=Inches(1)
			txBox = EG_slide.shapes.add_textbox(tl, tt, tw, th)
			tf = txBox.text_frame
			tf.text = 'forming bond '+FBL+' > ligand '+ LGL + ' > Leaving group (class) ' + EG
		# Create a go back button to the plot slide
			bl=Inches(17)
			bt=Inches(0)
			bw=Inches(1)
			bh=Inches(1)
		# Link this slide to example picture.
			rxn_click_act = example_reaction_fig.click_action
			rxn_click_act.target_slide = EG_slide
			cluster_click_act.action
			backbuttonEG = EG_slide.shapes.add_shape(MSO_SHAPE.ACTION_BUTTON_BACK_OR_PREVIOUS, bl, bt, bw, bh)
		# Create back button on new slide
			bbEG_click_act = backbuttonEG.click_action
			bbEG_click_act.target_slide = cluster_slide
			bbEG_click_act.action
			for egsub_index,EG_sub in enumerate(EG_sub_dirs):
				print('EG_sub',EG_sub)
			# Add ref for each example
				leftpin=coord_grid[egsub_index][0]
				toppin=coord_grid[egsub_index][1]
				EG_sub_path=EG_path+'/'+EG_sub
			# Get one figure as example
				file_paths=absoluteFilePaths(EG_sub_path)
				example_reaction_path=file_paths[0]
				e_rxn = EG_slide.shapes.add_picture(example_reaction_path, Inches(leftpin+0.2), Inches(toppin+0.2), height=Inches(1))
			# Add mechanism figure contain that existing group
				mechanism=False
				for cpath in file_paths:
					if 'mechanism' in cpath:
						mech_path=cpath

						mechanism=True
						break
				if mechanism:
					mech = EG_slide.shapes.add_picture(mech_path, Inches(12), Inches(5), height=Inches(4))
				# Add ref to that mechanism
					tl=Inches(12)
					tt=Inches(4)
					tw=Inches(6)
					th=Inches(1)
					txBox_mech_ref = EG_slide.shapes.add_textbox(tl, tt, tw, th)
					tf_mech_ref = txBox_mech_ref.text_frame
					mech_ref=mech_path.split('mechanism')[1]
					tf_mech_ref.text = mech_ref.replace('.png','')
			# Add text to cluster's slide
				tl=Inches(leftpin)
				tt=Inches(toppin)
				tw=Inches(1)
				th=Inches(1)
				txBox_EG_sub = EG_slide.shapes.add_textbox(tl, tt, tw, th)
				tf_EG_sub = txBox_EG_sub.text_frame
				tf_EG_sub.text = EG_sub
				tl=Inches(leftpin+1)
				tt=Inches(toppin)
				tw=Inches(1)
				th=Inches(1)
				txBox_ref = EG_slide.shapes.add_textbox(tl, tt, tw, th)
				tf_ref = txBox_ref.text_frame
	#			print('example_reaction_path',example_reaction_path)
				part_of_paths=example_reaction_path[-25:].split('\\')
				ID_file_name=part_of_paths[-1]
				print('file_name',ID_file_name)
				tf_ref.text = ID_file_name.replace('.png','')
#----------------------------------------------------------------------------------
prs = Presentation()
H_length=18
V_length=12
prs.slide_width = Inches(H_length)
prs.slide_height = Inches(V_length)
# Adding plot to the slide
blank_slide_layout = prs.slide_layouts[6]
plot_slide = prs.slides.add_slide(blank_slide_layout)
path='2_plot.png'
# Adding picture to the slide
plot = plot_slide.shapes.add_picture(path, Inches(0), Inches(0), height=Inches(12))
coor_data=pd.read_excel('2_cluster_coord.xlsx')
#print(coor_data)

size_H_max=coor_data.LG_C.max()+1
size_V_max=coor_data.FB_C.max()+1

rect_size_H=H_length*0.35/size_H_max
rect_size_V=V_length*0.967/size_V_max

# GO THROUGH EACH CLUSTER
for index,row in coor_data.iterrows():
	FBL=str(row['FB_L'])
	LGL=str(row['LG_L'])
	if LGL=='nan': continue
	FBC=row['FB_C']
	LGC=row['LG_C']
	rl=Inches(5.5+LGC*rect_size_H)
	rt=Inches((size_V_max-1-FBC)*rect_size_V)
	rw=Inches(rect_size_H)
	rh=Inches(rect_size_V)
	rect = plot_slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rl, rt, rw, rh)
	cluster_slide = prs.slides.add_slide(blank_slide_layout)#-----------------------------------------------------
	tl=Inches(1)
	tt=Inches(0)
	tw=Inches(15)
	th=Inches(1)
	# Add text to cluster's slide
	txBox = cluster_slide.shapes.add_textbox(tl, tt, tw, th)
	tf = txBox.text_frame
	tf.text = 'forming bond '+FBL+' > ligand '+ LGL
	different_scafold=False
	for special_bond in ['Csp3-Csp3','Csp2-Csp2','Csp1-Csp1','Csp1=Csp1']:
		if special_bond in FBL:
			different_scafold=True
			break
	# Add scaffold information (if have)
	if not different_scafold:
		tl=Inches(2)
		tt=Inches(0.5)
		tw=Inches(15)
		th=Inches(1)
		txBox1 = cluster_slide.shapes.add_textbox(tl, tt, tw, th)
		tf = txBox1.text_frame
		tf.text = 'Scaffold (left)'
		tl=Inches(8)
		tt=Inches(0.5)
		tw=Inches(15)
		th=Inches(1)
		txBox2 = cluster_slide.shapes.add_textbox(tl, tt, tw, th)
		tf = txBox2.text_frame
		tf.text = 'Scaffold (right)'
	# Create a go back button to the plot slide
	bl=Inches(17)
	bt=Inches(0)
	bw=Inches(1)
	bh=Inches(1)
	backbutton = cluster_slide.shapes.add_shape(MSO_SHAPE.ACTION_BUTTON_BACK_OR_PREVIOUS, bl, bt, bw, bh)
	bb_click_act = backbutton.click_action
	bb_click_act.target_slide = plot_slide
	bb_click_act.action
	# Each cluster has link to slide
	cluster_click_act = rect.click_action
	cluster_click_act.target_slide = cluster_slide
	cluster_click_act.action
	# GO TO CLUSTER FOLDER and GET FILE LIST
	cluster_path='2_library/'+FBL+'_'+LGL
	cluster_path=cluster_path.replace(' ','')
	print('cluster_path',cluster_path)
	thedir=cluster_path
	cluster_sub_dirs=([ name for name in os.listdir(thedir) if os.path.isdir(os.path.join(thedir, name)) ])
	# Set up some grid
	coord_grid=([0,1],[0,2],[0,3],[0,4],[0,5],[0,6],[0,7],[0,8],[0,9],[0,10],[6,1],[6,2],[6,3],[6,4],[6,5],[6,6],[6,7],[6,8],[6,9],[6,10])
	# GO TO EACH LEAVING GROUP FOLDER
	# First case: two scaffold are the same
	if 'EG' in cluster_sub_dirs:
		scaffold_path=cluster_path+'/EG'
		EG_slide()
	# TWO DIFFERENT SCAFFOLD CASES
	else:
		scaffold_path=cluster_path+'/EG1'
		EG_slide(offset=0)
		scaffold_path=cluster_path+'/EG2'
		EG_slide(offset=6)
prs.save('result.pptx')
